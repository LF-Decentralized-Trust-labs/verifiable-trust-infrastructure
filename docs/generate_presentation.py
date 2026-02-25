#!/usr/bin/env python3
"""Generate VTA Architecture Overview PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT_TEAL = RGBColor(0x00, 0x96, 0x88)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Calibri"


def set_font(run, size=14, bold=False, color=DARK_TEXT, name=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_title_bar(slide, text, top=Inches(0), height=Inches(0.9)):
    """Add a dark blue title bar at the top of a slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_font(run, size=28, bold=True, color=WHITE)
    return bar


def add_bullet_frame(slide, left, top, width, height):
    """Add a text box and return its text_frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_bullet(tf, text, level=0, size=14, bold=False, color=DARK_TEXT, spacing_before=4):
    """Add a bullet paragraph to a text frame."""
    p = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(spacing_before)
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return p


def add_section_header(tf, text, size=18, spacing=8):
    """Add a bold section header bullet."""
    return add_bullet(tf, text, level=0, size=size, bold=True, color=ACCENT_BLUE, spacing_before=spacing)


def add_box(slide, left, top, width, height, fill_color=LIGHT_GRAY, border_color=ACCENT_BLUE):
    """Add a rounded rectangle box."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape


def make_component_box(slide, left, top, width, height, title, bullets, fill=LIGHT_GRAY, border=ACCENT_BLUE):
    """Create a labeled box with bullet text inside."""
    box = add_box(slide, left, top, width, height, fill, border)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, size=14, bold=True, color=DARK_BLUE)

    for b in bullets:
        p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"  {b}"
        set_font(run, size=10, color=DARK_TEXT)

    return box


# ── Build the presentation ──────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # blank layout

# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 1 — Title                                               ║
# ╚══════════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(blank)
# background fill
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = DARK_BLUE

# Accent band
band = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.0), SLIDE_W, Inches(0.08)
)
band.fill.solid()
band.fill.fore_color.rgb = ACCENT_TEAL
band.line.fill.background()

# Title
tb = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.4))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Verifiable Trust Agent (VTA)"
set_font(run, size=44, bold=True, color=WHITE)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "Architecture Overview"
set_font(run2, size=36, bold=False, color=RGBColor(0xBB, 0xDE, 0xFB))

# Subtitle
tb2 = slide.shapes.add_textbox(Inches(1), Inches(3.4), Inches(11), Inches(1.2))
tf2 = tb2.text_frame
tf2.word_wrap = True
p3 = tf2.paragraphs[0]
p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
run3.text = "First Person Network  |  February 2026"
set_font(run3, size=20, color=RGBColor(0x90, 0xCA, 0xF9))

p4 = tf2.add_paragraph()
p4.alignment = PP_ALIGN.CENTER
p4.space_before = Pt(8)
run4 = p4.add_run()
run4.text = "vtc-vta-rs workspace: VTA Service \u00b7 VTC Service \u00b7 PNM CLI \u00b7 CNM CLI"
set_font(run4, size=14, color=RGBColor(0x78, 0x9E, 0xC2))


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 2 — System Overview                                     ║
# ╚══════════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(blank)
add_title_bar(slide, "System Overview")

tf = add_bullet_frame(slide, Inches(0.6), Inches(1.1), Inches(5.8), Inches(5.8))

# What is FPN
add_section_header(tf, "First Person Network (FPN)", spacing=4)
add_bullet(tf, "A decentralized identity & trust framework where individuals and")
add_bullet(tf, "communities own their cryptographic keys and credentials", level=0)
add_bullet(tf, "Anchored on DID-based identity (did:webvh, did:key)", level=1, size=12, color=MEDIUM_GRAY)

add_section_header(tf, "Verifiable Trust Agent (VTA)")
add_bullet(tf, "Personal cryptographic agent \u2014 manages keys, DIDs, credentials")
add_bullet(tf, "Runs as a local service; communicates via REST API and DIDComm v2")
add_bullet(tf, "BIP-32 hierarchical deterministic key derivation from a BIP-39 mnemonic")
add_bullet(tf, "Organizes keys into named contexts (e.g., 'vta', 'mediator', 'community')")

add_section_header(tf, "Verifiable Trust Community (VTC)")
add_bullet(tf, "Community-level service \u2014 manages membership, ACLs")
add_bullet(tf, "No local key management; receives keys from a VTA")
add_bullet(tf, "Provides authentication and authorization for community members")

# Right side: concept diagram as text boxes
make_component_box(slide, Inches(7.0), Inches(1.2), Inches(2.4), Inches(1.2),
    "Person", ["Owns a VTA", "Manages personal keys", "Joins communities"], border=ACCENT_TEAL)
make_component_box(slide, Inches(9.8), Inches(1.2), Inches(2.8), Inches(1.2),
    "Community", ["Runs a VTC", "Defines membership ACL", "Receives keys from VTAs"], border=ACCENT_BLUE)

# Arrow representation
arr = slide.shapes.add_textbox(Inches(7.2), Inches(2.6), Inches(5.2), Inches(0.5))
atf = arr.text_frame
p = atf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "\u2500\u2500\u2500\u2500  DIDComm v2 / REST  \u2500\u2500\u2500\u2500\u25b6"
set_font(run, size=12, color=ACCENT_BLUE)

make_component_box(slide, Inches(7.0), Inches(3.3), Inches(5.6), Inches(1.5),
    "WebVH Server", [
        "Hosts did:webvh documents",
        "REST API + DIDComm transport",
        "DID lifecycle: allocate \u2192 publish \u2192 update \u2192 deactivate"
    ], border=RGBColor(0x8B, 0xC3, 0x4A))

make_component_box(slide, Inches(7.0), Inches(5.2), Inches(5.6), Inches(1.5),
    "DIDComm Mediator (ATM)", [
        "Routes encrypted DIDComm messages between agents",
        "WebSocket-based live delivery",
        "Message pickup & forwarding"
    ], border=RGBColor(0xFF, 0x98, 0x00))


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 3 — Component Overview                                  ║
# ╚══════════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(blank)
add_title_bar(slide, "Component Overview \u2014 Workspace at a Glance")

components = [
    ("VTA Service", "vta-service", [
        "Personal cryptographic agent",
        "REST API + DIDComm messaging",
        "BIP-32 key derivation, contexts, ACL",
        "fjall embedded KV store",
        "Default port: 8100"
    ]),
    ("VTC Service", "vtc-service", [
        "Community trust service",
        "REST API + DIDComm messaging",
        "No local key management",
        "Membership ACL only",
        "Default port: 8200"
    ]),
    ("PNM CLI", "pnm-cli", [
        "Personal Network Manager",
        "Single VTA management",
        "Interactive setup wizard",
        "WebVH DID operations",
        "Keys, contexts, ACL commands"
    ]),
    ("CNM CLI", "cnm-cli", [
        "Community Network Manager",
        "Multi-community support",
        "Community add/use/remove",
        "Cross-VTA session bootstrap",
        "Shared SDK with PNM"
    ]),
    ("VTA SDK", "vta-sdk", [
        "Shared client library",
        "REST + DIDComm transport",
        "Auth session management",
        "Used by PNM and CNM CLIs",
        "Typed API wrappers"
    ]),
]

for i, (title, crate, bullets) in enumerate(components):
    x = Inches(0.4 + i * 2.55)
    make_component_box(slide, x, Inches(1.2), Inches(2.35), Inches(2.6), title, bullets)
    # crate label
    lb = slide.shapes.add_textbox(x, Inches(3.85), Inches(2.35), Inches(0.4))
    ltf = lb.text_frame
    p = ltf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = crate
    set_font(run, size=10, color=MEDIUM_GRAY)

# Shared components note
tf = add_bullet_frame(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(2.5))
add_section_header(tf, "Shared Infrastructure", spacing=4)
add_bullet(tf, "affinidi-tdk 0.4 \u2014 DIDComm messaging, secrets resolver, DID crypto")
add_bullet(tf, "affinidi-did-resolver-cache-sdk \u2014 DID resolution with in-memory cache")
add_bullet(tf, "fjall \u2014 Embedded LSM-tree key-value store (keys, sessions, ACL, contexts, webvh keyspaces)")
add_bullet(tf, "jsonwebtoken 10 with EdDSA (Ed25519) \u2014 JWT-based REST authentication")


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 4 — VTA Service Deep Dive                               ║
# ╚══════════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(blank)
add_title_bar(slide, "VTA Service \u2014 Deep Dive")

# Left column: data model
tf = add_bullet_frame(slide, Inches(0.5), Inches(1.1), Inches(6.0), Inches(5.8))

add_section_header(tf, "Key Management", spacing=4)
add_bullet(tf, "BIP-39 mnemonic \u2192 BIP-32 HD key derivation tree")
add_bullet(tf, "Ed25519 signing keys + X25519 key-agreement keys")
add_bullet(tf, "Keys organized into named contexts with derivation path counters")
add_bullet(tf, "Path format:  m/purpose'/context'/index'", size=12, color=MEDIUM_GRAY)

add_section_header(tf, "Contexts")
add_bullet(tf, "Named groupings of keys (e.g., 'vta', 'mediator', 'community')")
add_bullet(tf, "Each context has a description and auto-incrementing key index")
add_bullet(tf, "Context-scoped access control for admin delegation")

add_section_header(tf, "ACL & Credentials")
add_bullet(tf, "DID-based ACL entries with roles: Admin, Initiator, Application")
add_bullet(tf, "Admin can be super (all contexts) or scoped to specific contexts")
add_bullet(tf, "Credential bundles issued for automated auth (base64url-encoded)")

add_section_header(tf, "Storage \u2014 fjall Keyspaces")
add_bullet(tf, "keys \u2014 key records, path counters, seed metadata")
add_bullet(tf, "sessions \u2014 auth sessions + refresh token reverse index")
add_bullet(tf, "acl \u2014 DID-to-role/context mappings")
add_bullet(tf, "contexts \u2014 context records + global counter")
add_bullet(tf, "webvh \u2014 server and DID records (feature-gated)")

# Right column: thread model
make_component_box(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(1.6),
    "Thread: vta-rest", [
        "Axum HTTP server on dedicated single-threaded Tokio runtime",
        "All REST routes: health, auth, keys, contexts, ACL, config, webvh",
        "Holds AppState with all keyspace handles and DIDComm bridge",
        "Graceful shutdown via watch::channel"
    ])

make_component_box(slide, Inches(7.0), Inches(3.1), Inches(5.8), Inches(1.6),
    "Thread: vta-didcomm", [
        "ATM mediator connection over WebSocket",
        "Loads signing + key-agreement secrets into ThreadedSecretsResolver",
        "Dispatches inbound messages to protocol handlers",
        "Publishes DIDCommBridge for outbound messages from REST thread"
    ])

make_component_box(slide, Inches(7.0), Inches(5.0), Inches(5.8), Inches(1.3),
    "Thread: vta-storage", [
        "Periodic session cleanup (expired challenges + sessions)",
        "Flushes fjall to disk on shutdown (always last thread to join)",
        "Configurable cleanup interval"
    ])


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 5 — VTC Service                                         ║
# ╚══════════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(blank)
add_title_bar(slide, "VTC Service \u2014 Community Trust Service")

tf = add_bullet_frame(slide, Inches(0.5), Inches(1.1), Inches(6.0), Inches(5.8))

add_section_header(tf, "Purpose", spacing=4)
add_bullet(tf, "Manages community membership and access control")
add_bullet(tf, "Acts as the trust anchor for a community")
add_bullet(tf, "Authenticates members via DID challenge-response")

add_section_header(tf, "Key Differences from VTA")
add_bullet(tf, "No BIP-32 key derivation \u2014 receives 64 raw bytes from VTA")
add_bullet(tf, "  (32 bytes Ed25519 + 32 bytes X25519)", level=0, size=12, color=MEDIUM_GRAY)
add_bullet(tf, "No contexts or key management routes")
add_bullet(tf, "Only 2 keyspaces: sessions and acl (vs. 5 in VTA)")
add_bullet(tf, "SecretStore trait instead of SeedStore (get/set/delete)")

add_section_header(tf, "Routes")
add_bullet(tf, "GET /health")
add_bullet(tf, "POST /auth/challenge, /auth/, /auth/refresh, /auth/credentials")
add_bullet(tf, "GET|DELETE /auth/sessions, DELETE /auth/sessions/{id}")
add_bullet(tf, "GET|PATCH /config")
add_bullet(tf, "GET|POST /acl, GET|PATCH|DELETE /acl/{did}")

# Side-by-side comparison table as boxes
make_component_box(slide, Inches(7.0), Inches(1.2), Inches(2.8), Inches(3.0),
    "VTA", [
        "Port: 8100",
        "JWT Audience: VTA",
        "Env prefix: VTA_",
        "BIP-32 key derivation",
        "5 fjall keyspaces",
        "Keys + Contexts routes",
        "WebVH integration",
        "SeedStore trait",
        "Service type: VTARest"
    ], border=ACCENT_TEAL)

make_component_box(slide, Inches(10.2), Inches(1.2), Inches(2.8), Inches(3.0),
    "VTC", [
        "Port: 8200",
        "JWT Audience: VTC",
        "Env prefix: VTC_",
        "Raw key bytes from VTA",
        "2 fjall keyspaces",
        "ACL routes only",
        "No WebVH",
        "SecretStore trait",
        "Svc: VerifiableTrustCommunity"
    ], border=ACCENT_BLUE)

# VTC config note
make_component_box(slide, Inches(7.0), Inches(4.6), Inches(5.8), Inches(1.0),
    "VTC Configuration", [
        "vta_did \u2014 points to the community's VTA",
        "Same config.toml structure, same seed storage backends"
    ], border=MEDIUM_GRAY)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 6 — CLI Tools                                           ║
# ╚══════════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(blank)
add_title_bar(slide, "CLI Tools: PNM & CNM")

# PNM column
make_component_box(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(5.5),
    "PNM \u2014 Personal Network Manager", [
        "",
        "Manages a single personal VTA instance",
        "",
        "Commands:",
        "  health          \u2014  Service health + DID resolution + mediator trust-ping",
        "  auth login/logout/status  \u2014  Session management",
        "  config get/update",
        "  keys list/create/get/revoke/rename",
        "  keys secrets/seeds/rotate-seed",
        "  contexts list/get/create/update/delete/bootstrap",
        "  acl list/get/create/update/delete",
        "  auth-credential create",
        "  setup           \u2014  Interactive VTA setup wizard",
        "",
        "WebVH Commands (PNM only):",
        "  webvh add-server / list-servers / update-server / remove-server",
        "  webvh create-did / list-dids / get-did / delete-did",
    ], border=ACCENT_TEAL)

# CNM column
make_component_box(slide, Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.5),
    "CNM \u2014 Community Network Manager", [
        "",
        "Manages one or more community VTAs (multi-community)",
        "",
        "Additional Commands:",
        "  community list/use/add/remove/status/ping",
        "",
        "Key Differences from PNM:",
        "  \u2022  CnmConfig with communities map (named communities)",
        "  \u2022  Per-community keyring credentials",
        "  \u2022  community use  \u2014  switch active community",
        "  \u2022  community ping  \u2014  DIDComm trust-ping to community VTA",
        "  \u2022  Optional personal_vta pointer for session bootstrapping",
        "  \u2022  No webvh commands (managed via PNM/VTA directly)",
        "",
        "Shared Foundation:",
        "  Both CLIs use vta-sdk and vta-cli-common crate",
        "  Same auth flow, same REST/DIDComm transport layer",
    ], border=ACCENT_BLUE)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 7 — WebVH Server Integration                            ║
# ╚══════════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(blank)
add_title_bar(slide, "WebVH Server Integration")

tf = add_bullet_frame(slide, Inches(0.5), Inches(1.1), Inches(6.0), Inches(2.5))
add_section_header(tf, "What is did:webvh?", spacing=4)
add_bullet(tf, "Verifiable History DID method \u2014 web-hosted DID documents with")
add_bullet(tf, "cryptographic audit trail (did.jsonl log entries)")
add_bullet(tf, "Supports portability and pre-rotation keys", level=0, size=12, color=MEDIUM_GRAY)

add_section_header(tf, "Server Registration")
add_bullet(tf, "Servers stored in fjall webvh keyspace with ID, DID, and label")
add_bullet(tf, "VTA resolves server DID document to find WebVHHostingService endpoint")
add_bullet(tf, "Picks REST vs DIDComm transport based on available endpoints")

# REST transport box
make_component_box(slide, Inches(0.4), Inches(3.9), Inches(6.0), Inches(2.8),
    "REST Transport", [
        "",
        "POST /api/dids           \u2014  Allocate DID URI (returns did_url + mnemonic)",
        "PUT  /api/dids/{mnemonic} \u2014  Publish signed did.jsonl",
        "DELETE /api/dids/{mnemonic} \u2014  Remove a DID",
        "POST /api/dids/check      \u2014  Check path availability",
        "",
        "Auth: Bearer {access_token} when configured",
        "Content-Type: application/jsonl for publish",
    ], border=ACCENT_TEAL)

# DIDComm transport box
make_component_box(slide, Inches(6.8), Inches(1.2), Inches(6.0), Inches(2.3),
    "DIDComm Transport", [
        "",
        "Routes through ATM mediator via DIDCommBridge::send_and_wait()",
        "Protocol: https://affinidi.com/webvh/1.0/did/",
        "",
        "did/request  \u2192  did/offer       (allocate URI)",
        "did/publish  \u2192  did/confirm     (publish did.jsonl)",
        "did/delete   \u2192  did/delete-confirm",
        "Errors: did/problem-report    (30s timeout per operation)",
    ], border=ACCENT_BLUE)

# DID lifecycle box
make_component_box(slide, Inches(6.8), Inches(3.9), Inches(6.0), Inches(2.8),
    "DID Lifecycle via VTA", [
        "",
        "1. Derive Ed25519 + X25519 keys in a context",
        "2. Request DID URI from WebVH server",
        "3. Build DID document with verification methods",
        "   (authentication, assertionMethod, keyAgreement)",
        "4. Add service endpoints (mediator, REST)",
        "5. Sign and publish did.jsonl log entry",
        "6. Store DID record in fjall webvh keyspace",
        "7. Update/deactivate via signed log entries",
    ], border=RGBColor(0x8B, 0xC3, 0x4A))


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 8 — Communication Architecture                          ║
# ╚══════════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(blank)
add_title_bar(slide, "Communication Architecture")

# Main diagram using text boxes and shapes
# Row 1: CLI clients
make_component_box(slide, Inches(0.4), Inches(1.4), Inches(2.0), Inches(0.9),
    "PNM CLI", ["REST + DIDComm"], border=ACCENT_TEAL)
make_component_box(slide, Inches(0.4), Inches(2.6), Inches(2.0), Inches(0.9),
    "CNM CLI", ["REST + DIDComm"], border=ACCENT_BLUE)

# Arrow 1: CLIs -> Mediator
a1 = slide.shapes.add_textbox(Inches(2.5), Inches(1.8), Inches(1.8), Inches(0.5))
p = a1.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "\u2500\u2500 DIDComm \u2500\u25b6"
set_font(run, size=11, bold=True, color=ACCENT_BLUE)

# REST arrow
a1r = slide.shapes.add_textbox(Inches(2.5), Inches(2.7), Inches(1.8), Inches(0.5))
p = a1r.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "\u2500\u2500\u2500 REST \u2500\u2500\u25b6"
set_font(run, size=11, bold=True, color=ACCENT_TEAL)

# Mediator
make_component_box(slide, Inches(4.4), Inches(1.4), Inches(2.2), Inches(1.2),
    "ATM Mediator", ["WebSocket relay", "Message pickup"], border=RGBColor(0xFF, 0x98, 0x00))

# Arrow 2: Mediator -> VTA
a2 = slide.shapes.add_textbox(Inches(6.7), Inches(1.6), Inches(1.3), Inches(0.5))
p = a2.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "\u25c0\u2500\u2500\u2500\u25b6"
set_font(run, size=11, bold=True, color=ACCENT_BLUE)

# VTA Service
make_component_box(slide, Inches(8.1), Inches(1.2), Inches(2.4), Inches(1.6),
    "VTA Service", ["REST: 0.0.0.0:8100", "DIDComm: via ATM WS", "fjall KV store"],
    border=ACCENT_TEAL)

# Arrow 3: VTA -> Mediator -> WebVH
a3 = slide.shapes.add_textbox(Inches(10.6), Inches(1.6), Inches(1.0), Inches(0.5))
p = a3.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "\u25c0\u2500\u2500\u25b6"
set_font(run, size=11, bold=True, color=RGBColor(0x8B, 0xC3, 0x4A))

# WebVH
make_component_box(slide, Inches(11.0), Inches(1.2), Inches(2.0), Inches(1.6),
    "WebVH Server", ["REST API", "DIDComm via ATM"],
    border=RGBColor(0x8B, 0xC3, 0x4A))

# Direct REST path (CLI -> VTA)
a_rest = slide.shapes.add_textbox(Inches(2.5), Inches(3.5), Inches(8.0), Inches(0.4))
p = a_rest.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "CLI \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500  Direct REST (HTTPS)  \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u25b6 VTA"
set_font(run, size=12, bold=False, color=ACCENT_TEAL)

# Explanation text
tf = add_bullet_frame(slide, Inches(0.5), Inches(4.2), Inches(12), Inches(3.0))
add_section_header(tf, "Two Communication Paths", spacing=4)
add_bullet(tf, "DIDComm Path: End-to-end encrypted messages routed through ATM mediator (WebSocket)")
add_bullet(tf, "  \u2022 CLI packs message \u2192 mediator relays \u2192 VTA unpacks and processes", level=0, size=12, color=MEDIUM_GRAY)
add_bullet(tf, "  \u2022 VTA \u2192 mediator \u2192 WebVH for DID operations (same encrypted path)", level=0, size=12, color=MEDIUM_GRAY)
add_bullet(tf, "REST Path: Direct HTTPS to VTA service (JWT-authenticated after challenge-response)")
add_bullet(tf, "  \u2022 Simpler, lower latency for direct access scenarios", level=0, size=12, color=MEDIUM_GRAY)
add_bullet(tf, "VTA SDK: Abstracts transport selection \u2014 same API regardless of path chosen")


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 9 — Authentication                                      ║
# ╚══════════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(blank)
add_title_bar(slide, "Authentication")

# REST auth flow
make_component_box(slide, Inches(0.4), Inches(1.2), Inches(6.2), Inches(4.0),
    "REST Authentication \u2014 DID Challenge-Response", [
        "",
        "1. POST /auth/challenge { did }",
        "   \u2192 VTA creates ChallengeSent session (TTL: 300s)",
        "   \u2192 Returns { session_id, challenge }",
        "",
        "2. Client builds DIDComm v2 message:",
        "   type: atm/1.0/authenticate",
        "   body: { challenge, session_id }",
        "   from: client_did, to: vta_did",
        "   \u2192 pack_encrypted() produces ciphertext",
        "",
        "3. POST /auth/ { packed_message }",
        "   \u2192 VTA unpacks, verifies challenge + DID match",
        "   \u2192 Looks up ACL entry for sender DID",
        "   \u2192 Issues EdDSA JWT (15 min) + refresh token (24 hr)",
        "",
        "4. Authorization: Bearer {JWT} on subsequent requests",
        "",
        "5. POST /auth/refresh { refresh_token }",
        "   \u2192 Re-reads ACL (role/contexts may have changed)",
        "   \u2192 Issues new access token",
    ], border=ACCENT_TEAL)

# DIDComm auth
make_component_box(slide, Inches(6.8), Inches(1.2), Inches(6.0), Inches(2.2),
    "DIDComm Authentication", [
        "",
        "Inherent \u2014 no separate auth step needed",
        "Message::unpack() verifies sender signature + resolves DID",
        "ACL looked up by 'from' DID (fragment stripped)",
        "Same role/context rules as REST path",
        "Every DIDComm message is authenticated by design",
    ], border=ACCENT_BLUE)

# ACL roles
make_component_box(slide, Inches(6.8), Inches(3.8), Inches(6.0), Inches(2.8),
    "ACL Roles & Privilege Model", [
        "",
        "Admin  \u2014  Full CRUD on keys, ACL, contexts, config",
        "  Super Admin: allowed_contexts = [] (unrestricted)",
        "  Context Admin: scoped to named contexts only",
        "",
        "Initiator  \u2014  Can manage ACL entries, view all resources",
        "",
        "Application  \u2014  Read-only access to keys, contexts, config",
        "",
        "Escalation prevention:",
        "  Context admins cannot grant broader contexts than their own",
        "  Cannot create entries with empty allowed_contexts",
    ], border=RGBColor(0xFF, 0x98, 0x00))

# JWT claims
make_component_box(slide, Inches(0.4), Inches(5.5), Inches(6.2), Inches(1.2),
    "JWT Claims", [
        '{ aud: "VTA", sub: "did:...", session_id: "uuid",',
        '  role: "admin", contexts: ["vta", "mediator"], exp: ... }',
        "Signed with random Ed25519 key (PKCS8 DER / SPKI DER)"
    ], border=MEDIUM_GRAY)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 10 — Data Flow Example                                  ║
# ╚══════════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(blank)
add_title_bar(slide, "Data Flow \u2014 Create a did:webvh via DIDComm")

steps = [
    ("1", "PNM CLI", "User runs: pnm webvh create-did\nSelects server, context, and options", ACCENT_TEAL),
    ("2", "VTA SDK", "SDK packs DIDComm message:\nfrom: user_did, to: vta_did\nRoutes through mediator WebSocket", ACCENT_BLUE),
    ("3", "VTA Service", "vta-didcomm thread receives message\nDerives new Ed25519 + X25519 keys\nBuilds DID document with endpoints", ACCENT_TEAL),
    ("4", "VTA \u2192 WebVH", "VTA sends did/request via DIDComm\n(or REST) to WebVH server\nReceives did/offer with allocated URI", RGBColor(0x8B, 0xC3, 0x4A)),
    ("5", "VTA \u2192 WebVH", "VTA signs did.jsonl log entry\nSends did/publish to WebVH\nReceives did/confirm", RGBColor(0x8B, 0xC3, 0x4A)),
    ("6", "VTA Service", "Stores key records in fjall keys KS\nStores DID record in webvh KS\nPrepares response message", ACCENT_TEAL),
    ("7", "VTA \u2192 PNM", "VTA packs DIDComm response\nRoutes through mediator to CLI\nPNM displays new DID to user", ACCENT_BLUE),
]

for i, (num, label, desc, color) in enumerate(steps):
    y = Inches(1.15 + i * 0.85)
    # Step number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y + Inches(0.05), Inches(0.45), Inches(0.45))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    ctf = circ.text_frame
    ctf.margin_left = Inches(0)
    ctf.margin_right = Inches(0)
    ctf.margin_top = Inches(0)
    p = ctf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    set_font(run, size=16, bold=True, color=WHITE)

    # Label
    lb = slide.shapes.add_textbox(Inches(1.1), y, Inches(2.0), Inches(0.55))
    p = lb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = label
    set_font(run, size=13, bold=True, color=color)

    # Description
    db = slide.shapes.add_textbox(Inches(3.2), y - Inches(0.05), Inches(9.5), Inches(0.7))
    dtf = db.text_frame
    dtf.word_wrap = True
    for j, line in enumerate(desc.split("\n")):
        if j == 0:
            p = dtf.paragraphs[0]
        else:
            p = dtf.add_paragraph()
        run = p.add_run()
        run.text = line
        set_font(run, size=11, color=DARK_TEXT)


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 11 — Setup & Configuration                              ║
# ╚══════════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(blank)
add_title_bar(slide, "Setup & Configuration")

# Setup wizard steps
make_component_box(slide, Inches(0.4), Inches(1.2), Inches(6.2), Inches(5.6),
    "VTA Setup Wizard  (pnm setup)", [
        "",
        " 1.  Config file path (default: config.toml)",
        " 2.  VTA name (optional community name)",
        " 3.  Services: REST API, DIDComm Messaging (multi-select)",
        " 4.  Public URL, host, port (if REST enabled)",
        " 5.  Log level + format (text/json)",
        " 6.  Data directory (opens fjall store)",
        " 7.  Creates default 'vta' application context",
        " 8.  BIP-39 mnemonic: generate 24-word or import existing",
        " 9.  Seed storage backend selection",
        "10.  JWT signing key generation (random Ed25519)",
        "11.  DIDComm: use existing mediator / create did:webvh / skip",
        "12.  VTA DID: create did:webvh / enter existing / skip",
        "13.  Admin DID: generate did:key / create did:webvh / enter / skip",
        "14.  Bootstrap ACL with admin entry",
        "15.  Persist store + save config.toml",
        "",
        "Feature-gated: requires 'setup' feature (excludes dialoguer",
        "from production builds)",
    ], border=ACCENT_TEAL)

# Config structure
make_component_box(slide, Inches(6.8), Inches(1.2), Inches(6.0), Inches(3.2),
    "config.toml Structure", [
        "",
        "[general]     \u2014  name, log_level, log_format",
        "[services]    \u2014  rest (bool), didcomm (bool)",
        "[rest]        \u2014  public_url, host, port",
        "[store]       \u2014  data_dir",
        "[auth]        \u2014  jwt_signing_key, challenge_ttl, session_ttl,",
        "                 session_cleanup_interval",
        "[secrets]     \u2014  seed storage backend config",
        "[didcomm]     \u2014  mediator_did",
        "[did]         \u2014  vta_did",
        "",
        "Env var overrides: VTA_{SECTION}_{KEY} (e.g., VTA_REST_PORT)",
    ], border=ACCENT_BLUE)

# CLI setup
make_component_box(slide, Inches(6.8), Inches(4.8), Inches(6.0), Inches(1.8),
    "CLI Setup", [
        "",
        "PNM: pnm setup \u2014 prompts for credential + VTA URL",
        "     Stores auth credential in OS keyring",
        "",
        "CNM: cnm setup \u2014 interactive community add flow",
        "     Per-community keyring entries",
        "     Optional personal_vta pointer for bootstrap",
    ], border=RGBColor(0xFF, 0x98, 0x00))


# ╔══════════════════════════════════════════════════════════════════╗
# ║  SLIDE 12 — Infrastructure & Dependencies                      ║
# ╚══════════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(blank)
add_title_bar(slide, "Infrastructure & Dependencies")

# Seed backends
make_component_box(slide, Inches(0.4), Inches(1.2), Inches(4.0), Inches(3.0),
    "Seed Storage Backends (by priority)", [
        "",
        "1. AWS Secrets Manager  (feature: aws-secrets)",
        "2. GCP Secret Manager   (feature: gcp-secrets)",
        "3. Azure Key Vault      (feature: azure-secrets)",
        "4. Config file          (feature: config-seed)",
        "5. OS Keyring           (feature: keyring)",
        "6. Plaintext file       (always compiled)",
        "",
        "All store hex-encoded seed/secret",
        "Plaintext: NOT recommended (warns on use)",
    ], border=RGBColor(0xFF, 0x98, 0x00))

# Cryptography
make_component_box(slide, Inches(4.7), Inches(1.2), Inches(4.0), Inches(3.0),
    "Cryptography", [
        "",
        "Ed25519  \u2014  Digital signatures (signing keys)",
        "X25519   \u2014  Key agreement (DIDComm encryption)",
        "BIP-32   \u2014  HD key derivation (via ed25519-dalek-bip32)",
        "BIP-39   \u2014  Mnemonic seed phrase (24 words)",
        "EdDSA    \u2014  JWT signing (jsonwebtoken 10)",
        "PKCS8 v1 \u2014  Ed25519 private key DER encoding",
        "SPKI     \u2014  Ed25519 public key DER encoding",
        "",
        "DIDComm v2 envelope encryption via affinidi-tdk",
    ], border=ACCENT_TEAL)

# External services
make_component_box(slide, Inches(9.0), Inches(1.2), Inches(4.0), Inches(3.0),
    "External Services", [
        "",
        "ATM Mediator \u2014  DIDComm message routing",
        "  WebSocket-based, live delivery + pickup",
        "",
        "DID Resolver Cache \u2014  DID document resolution",
        "  affinidi-did-resolver-cache-sdk",
        "  In-memory cache with TTL",
        "",
        "WebVH Server \u2014  did:webvh hosting",
        "  REST + DIDComm transport",
    ], border=ACCENT_BLUE)

# Key dependencies
make_component_box(slide, Inches(0.4), Inches(4.5), Inches(12.6), Inches(2.5),
    "Key Rust Dependencies", [
        "",
        "affinidi-tdk 0.4             \u2014  DIDComm messaging, secrets resolver, DID crypto, ATM client",
        "affinidi-did-resolver-cache-sdk  \u2014  DID resolution (not re-exported by tdk, explicit dependency)",
        "axum 0.8 + axum-extra 0.12   \u2014  HTTP framework (typed headers, Authorization<Bearer>)",
        "fjall                         \u2014  Embedded LSM-tree KV store (keys, sessions, acl, contexts, webvh keyspaces)",
        "jsonwebtoken 10 (rust_crypto) \u2014  EdDSA JWT creation + validation",
        "ed25519-dalek-bip32           \u2014  BIP-32 hierarchical deterministic key derivation",
        "azure_security_keyvault_secrets / azure_identity  \u2014  Azure Key Vault integration",
        "dialoguer (setup feature)     \u2014  Interactive terminal prompts for setup wizard",
    ], border=MEDIUM_GRAY)


# ── Save ─────────────────────────────────────────────────────────────
output_path = "docs/VTA_Service_Overview.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"  {len(prs.slides)} slides generated")
